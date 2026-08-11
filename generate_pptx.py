import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_path):
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # blank layout

    # Color Palette Definitions
    COLOR_BG = RGBColor(15, 23, 42)       # Slate 900 #0F172A
    COLOR_CARD_BG = RGBColor(30, 41, 59)  # Slate 800 #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700 #334155
    COLOR_CYAN = RGBColor(14, 165, 233)   # Sky 500 #0EA5E9
    COLOR_PURPLE = RGBColor(139, 92, 246) # Purple 500 #8B5CF6
    COLOR_WHITE = RGBColor(248, 250, 252) # Slate 50 #F8FAFC
    COLOR_GRAY = RGBColor(148, 163, 184)  # Slate 400 #94A3B8
    COLOR_RED = RGBColor(239, 68, 68)     # Red 500 #EF4444
    COLOR_GREEN = RGBColor(16, 185, 129)  # Emerald 500 #10B981
    COLOR_AMBER = RGBColor(245, 158, 11)  # Amber 500 #F59E0B

    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="EVENT-CONNECT — SERVERLESS TICKETING HUB"):
        # Category Tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    # Decorative background accent card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(11.733), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1.5)

    # Title box
    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.8))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ Event-Connect"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p2 = tf.add_paragraph()
    p2.text = "Universal Multi-API Event Manager & Ticketing Integration Hub"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    # Subtitle / Description
    sbox = slide1.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(10.9), Inches(1.2))
    stf = sbox.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "A high-performance AWS serverless architecture connecting multi-region event APIs,\nenforcing zero overbooking, and delivering $0 running cost while idle."
    sp.font.size = Pt(16)
    sp.font.color.rgb = COLOR_GRAY

    # Footer metadata badges
    fbox = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.8))
    ftf = fbox.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "📍 Live App: https://geekkwame.github.io/event-registration-system-sam/   |   ☁️ AWS us-west-1 & us-east-1"
    fp.font.size = Pt(14)
    fp.font.color.rgb = COLOR_PURPLE

    # -------------------------------------------------------------
    # SLIDE 2: The Problem
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "The Pitfalls of Traditional Event Registration Systems")

    problems = [
        ("🚨 1. Overbooking & Venue Violations", 
         "Standard online forms (like MS Forms or Google Forms) lack live capacity enforcement.\n\n"
         "• Registration forms remain open even when seats are full.\n"
         "• Venue capacity thresholds are breached, forcing manual rejections."),
        
        ("⚠️ 2. Duplicate Registrations", 
         "Without immediate validation or pass issuance, attendees resubmit forms multiple times.\n\n"
         "• Inflated headcount metrics in spreadsheets.\n"
         "• Dirty data requiring hours of manual email cross-referencing."),
        
        ("🌐 3. Multi-System Fragmentation", 
         "Organizations operating across multiple departments or AWS regions face disconnected APIs.\n\n"
         "• No single portal to view partner events.\n"
         "• Incompatible payload schemas and CORS cross-origin errors.")
    ]

    left_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]
    for i, (p_title, p_desc) in enumerate(problems):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[i], Inches(1.8), Inches(3.733), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_RED
        card.line.width = Pt(1.5)

        tbox = slide2.shapes.add_textbox(left_positions[i] + Inches(0.2), Inches(2.0), Inches(3.333), Inches(4.4))
        tf = tbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + p_desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 3: The Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "The Event-Connect Solution & Core Value Pillars")

    solutions = [
        ("🔗 Universal Multi-API Hub", 
         "Connect to any AWS SAM API Gateway endpoint across regions instantly using Quick-Connect presets.\n\n"
         "• Dynamic endpoint switching.\n"
         "• Automated stage path normalization.\n"
         "• Full CORS preflight compatibility."),
        
        ("🛡️ Automated Capacity Control", 
         "Backend microservices validate DynamoDB seat counts atomically before saving sign-ups.\n\n"
         "• Hard capacity caps block overbooking.\n"
         "• Live color-coded capacity badges.\n"
         "• One-click instant seat restoration."),
        
        ("💰 $0 Running Cost When Idle", 
         "Built strictly with AWS Serverless components to eliminate fixed monthly server costs.\n\n"
         "• Pay per request execution only.\n"
         "• Free tier eligible across all AWS services.\n"
         "• Zero server management overhead.")
    ]

    for i, (s_title, s_desc) in enumerate(solutions):
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[i], Inches(1.8), Inches(3.733), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_GREEN
        card.line.width = Pt(1.5)

        tbox = slide3.shapes.add_textbox(left_positions[i] + Inches(0.2), Inches(2.0), Inches(3.333), Inches(4.4))
        tf = tbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + s_desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "End-to-End AWS Serverless Architecture")

    arch_steps = [
        ("1. Web UI Dashboard", "GitHub Pages (HTML5/CSS3)\nGlassmorphism design, async fetch, API preset state management."),
        ("2. AWS API Gateway", "Secure REST API routing\nCORS preflight handling, stage normalization (/dev, /Prod)."),
        ("3. Python 3.12 Lambdas", "Microservices:\n• list_events.py\n• register.py\n• get_registrations.py\n• cancel_registration.py"),
        ("4. DynamoDB & SNS", "NoSQL Data Store:\n• Events & Registrations tables\n• EmailIndex GSI for fast lookup\n• SNS automated email alerts")
    ]

    step_lefts = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
    for i, (a_title, a_desc) in enumerate(arch_steps):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, step_lefts[i], Inches(2.0), Inches(2.733), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CYAN
        card.line.width = Pt(1.5)

        tbox = slide4.shapes.add_textbox(step_lefts[i] + Inches(0.15), Inches(2.2), Inches(2.433), Inches(4.1))
        tf = tbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = a_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + a_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 5: Multi-Region API Hub
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Universal Multi-Region API Integration")

    # Primary Hub Card
    card_hub = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_hub.fill.solid()
    card_hub.fill.fore_color.rgb = COLOR_CARD_BG
    card_hub.line.color.rgb = COLOR_CYAN
    card_hub.line.width = Pt(2.0)

    tbox = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🟢 Primary Hub API (AWS us-west-1)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "\nEndpoint: https://kems8drwn6.execute-api.us-west-1.amazonaws.com/Prod\n\n" \
              "• Full administrative & attendee capabilities.\n" \
              "• Production AWS SAM stack deployment.\n" \
              "• Real-time DynamoDB table synchronization.\n" \
              "• CloudWatch alarm monitoring at 5% error threshold."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_GRAY

    # Partner APIs Card
    card_part = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card_part.fill.solid()
    card_part.fill.fore_color.rgb = COLOR_CARD_BG
    card_part.line.color.rgb = COLOR_PURPLE
    card_part.line.width = Pt(2.0)

    tbox2 = slide5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf2 = tbox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "🟣 Partner Integration APIs (AWS us-east-1)"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE

    p4 = tf2.add_paragraph()
    p4.text = "\nEndpoints:\n" \
              "• Gloria's API: https://djabididt6.execute-api.us-east-1.amazonaws.com\n" \
              "• Dawuni's API: https://mmrq6ebalh.execute-api.us-east-1.amazonaws.com\n\n" \
              "• Cross-account regional event aggregation.\n" \
              "• Dynamic stage path handling.\n" \
              "• Unified event listing & ticketing in one dashboard."
    p4.font.size = Pt(14)
    p4.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 6: Attendee Lifecycle
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "Attendee Experience & Ticket Lifecycle")

    steps = [
        ("1. Event Discovery", "Browse events with live capacity progress bars and status indicators (Open, Near Full, Full)."),
        ("2. Smart Sign-Up", "Register with instant duplicate prevention, email format validation, and atomic seat locking."),
        ("3. Digital Pass Lookup", "Query active tickets anytime by email via DynamoDB EmailIndex Global Secondary Index (GSI)."),
        ("4. One-Click Cancel", "Cancel ticket instantly to remove registration and restore event capacity back to the public pool.")
    ]

    for i, (st_title, st_desc) in enumerate(steps):
        left_pos = Inches(0.8) + Inches(i * 3.0)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), Inches(2.7), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CYAN
        card.line.width = Pt(1.5)

        tbox = slide6.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.2), Inches(2.4), Inches(4.1))
        tf = tbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = st_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + st_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 7: Reliability & DevOps
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "DevOps, CI/CD Pipeline & Quality Assurance")

    # CI/CD Card
    card_cicd = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_cicd.fill.solid()
    card_cicd.fill.fore_color.rgb = COLOR_CARD_BG
    card_cicd.line.color.rgb = COLOR_CYAN
    card_cicd.line.width = Pt(1.5)

    tbox = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 GitHub Actions CI/CD Pipeline"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "\n• Automated Pull Request Checks: Executes unit tests automatically before merging.\n" \
              "• Continuous Deployment: Automatic build & deployment (`sam deploy`) to AWS on `main` push.\n" \
              "• GitHub Pages Automation: Auto-deploys static frontend updates on push."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_GRAY

    # Pytest & SAM Card
    card_test = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card_test.fill.solid()
    card_test.fill.fore_color.rgb = COLOR_CARD_BG
    card_test.line.color.rgb = COLOR_GREEN
    card_test.line.width = Pt(1.5)

    tbox2 = slide7.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf2 = tbox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "🧪 Pytest Suite & Infrastructure as Code"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE

    p4 = tf2.add_paragraph()
    p4.text = "\n• 9/9 Automated Unit Tests Passing: Powered by `moto` to mock AWS DynamoDB & SNS locally at $0 cost.\n" \
              "• AWS SAM IaC (`template.yaml`): Standardized CloudFormation definitions for all serverless resources.\n" \
              "• CloudWatch Monitoring: Error alarms configured for high availability."
    p4.font.size = Pt(14)
    p4.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------
    # SLIDE 8: Live System Demo Blueprint
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "Live Interactive System Demonstration")

    demo_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    demo_card.fill.solid()
    demo_card.fill.fore_color.rgb = COLOR_CARD_BG
    demo_card.line.color.rgb = COLOR_PURPLE
    demo_card.line.width = Pt(2.0)

    dtbox = slide8.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.4))
    dtf = dtbox.text_frame
    dtf.word_wrap = True

    dp = dtf.paragraphs[0]
    dp.text = "🖥️ Live Application Walkthrough Plan"
    dp.font.size = Pt(22)
    dp.font.bold = True
    dp.font.color.rgb = COLOR_CYAN

    dp2 = dtf.add_paragraph()
    dp2.text = "\n1. Quick-Connect API Presets: Switch between AWS us-west-1 Primary Hub & us-east-1 Partner APIs." \
               "\n2. Real-Time Capacity Tracking: Observe live event progress bars and capacity status badges." \
               "\n3. Attendee Registration Flow: Register attendee & receive instant Registration ID pass." \
               "\n4. Email Pass Lookup (GSI): Query attendee tickets dynamically from DynamoDB." \
               "\n5. Ticket Cancellation & Restoration: Cancel pass to verify instant seat capacity restoration."
    dp2.font.size = Pt(16)
    dp2.font.color.rgb = COLOR_WHITE

    dp3 = dtf.add_paragraph()
    dp3.text = "\n👉 Live Web App URL: https://geekkwame.github.io/event-registration-system-sam/"
    dp3.font.size = Pt(16)
    dp3.font.bold = True
    dp3.font.color.rgb = COLOR_GREEN

    # -------------------------------------------------------------
    # SLIDE 9: Impact & Roadmap
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "Summary of Impact & Future Roadmap")

    # Impact Card
    card_imp = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_imp.fill.solid()
    card_imp.fill.fore_color.rgb = COLOR_CARD_BG
    card_imp.line.color.rgb = COLOR_GREEN
    card_imp.line.width = Pt(1.5)

    tbox = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌟 Key Business Impact"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "\n• 100% Elimination of Overbooking: Hard backend capacity locks.\n" \
              "• 100% Duplicate Prevention: Enforced GSI unique email validation.\n" \
              "• $0 Monthly Idle Overhead: Pure pay-per-request serverless stack.\n" \
              "• Cross-Region Integration: Aggregate partner APIs seamlessly."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_GRAY

    # Roadmap Card
    card_road = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card_road.fill.solid()
    card_road.fill.fore_color.rgb = COLOR_CARD_BG
    card_road.line.color.rgb = COLOR_CYAN
    card_road.line.width = Pt(1.5)

    tbox2 = slide9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf2 = tbox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "🚀 Future System Roadmap"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE

    p4 = tf2.add_paragraph()
    p4.text = "\n• QR Code Ticket Check-in: Mobile camera scanning for on-site event staff.\n" \
              "• Webhook Integrations: Google Calendar & Outlook calendar sync.\n" \
              "• Multi-Tenant RBAC: Role-based access control for event organizers.\n\n" \
              "❓ Questions & Answers (Q&A)"
    p4.font.size = Pt(14)
    p4.font.color.rgb = COLOR_GRAY

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    out_file = sys.argv[1] if len(sys.argv) > 1 else "presentation.pptx"
    build_presentation(out_file)
